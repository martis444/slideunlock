import logging
from io import BytesIO

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

log = logging.getLogger(__name__)

_MSO = MSO_AUTO_SHAPE_TYPE

_AUTO_SHAPE_MAP = {
    "rectangle":            _MSO.RECTANGLE,
    "rounded_rect":         _MSO.ROUNDED_RECTANGLE,
    "oval":                 _MSO.OVAL,
    "triangle":             _MSO.ISOSCELES_TRIANGLE,
    "diamond":              _MSO.DIAMOND,
    "pentagon":             _MSO.PENTAGON,
    "hexagon":              _MSO.HEXAGON,
    "arrow_right":          _MSO.RIGHT_ARROW,
    "arrow_left":           _MSO.LEFT_ARROW,
    "arrow_double":         _MSO.LEFT_RIGHT_ARROW,
    "callout_rect":         _MSO.RECTANGULAR_CALLOUT,
    "callout_rounded_rect": _MSO.ROUNDED_RECTANGULAR_CALLOUT,
}

_CONNECTOR_MAP = {
    "straight": MSO_CONNECTOR_TYPE.STRAIGHT,
    "elbow":    MSO_CONNECTOR_TYPE.ELBOW,
    "curved":   MSO_CONNECTOR_TYPE.CURVE,
}

_VALIGN_MAP = {
    "top":    MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

_ALIGN_MAP = {
    "left":    PP_ALIGN.LEFT,
    "center":  PP_ALIGN.CENTER,
    "right":   PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

_CONNECTOR_TYPES = frozenset(["line", "connector"])


# ── per-shape helpers ────────────────────────────────────────────────────────

def _apply_rotation(shape, spec: dict) -> None:
    rot = spec.get("rot", 0)
    if not rot:
        return
    try:
        shape.element.spPr.xfrm.set("rot", str(rot))
    except AttributeError:
        log.debug("rotation skipped for shape id=%s (no xfrm)", spec.get("id"))


def _apply_fill(shape, spec: dict) -> None:
    fill_type = spec.get("fill_type", "none")

    if fill_type == "solid":
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(
            spec.get("fill_hex", "FFFFFF")
        )
        opacity = spec.get("fill_opacity", 1.0)
        if opacity < 1.0:
            solid_fill = shape.element.spPr.find(qn("a:solidFill"))
            if solid_fill is not None and len(solid_fill):
                alpha = etree.SubElement(solid_fill[0], qn("a:alpha"))
                alpha.set("val", str(int(opacity * 100000)))

    elif fill_type == "gradient":
        spPr = shape.element.spPr
        for tag in (
            qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
            qn("a:blipFill"), qn("a:pattFill"), qn("a:grpFill"),
        ):
            for old in spPr.findall(tag):
                spPr.remove(old)

        grad_fill = etree.SubElement(spPr, qn("a:gradFill"))
        gs_lst    = etree.SubElement(grad_fill, qn("a:gsLst"))

        for stop in spec.get("gradient_stops", []):
            gs   = etree.SubElement(gs_lst, qn("a:gs"))
            gs.set("pos", str(int(stop["pos"] * 100000)))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", stop["hex"])
            a_el = etree.SubElement(srgb, qn("a:alpha"))
            a_el.set("val", str(int(stop.get("opacity", 1.0) * 100000)))

        lin = etree.SubElement(grad_fill, qn("a:lin"))
        lin.set("ang",    str(int(spec.get("gradient_angle_deg", 90) * 60000)))
        lin.set("scaled", "0")

    else:  # "none"
        shape.fill.background()


def _apply_line(shape, spec: dict) -> None:
    line_width = spec.get("line_width_pt", 0.0)
    stype      = spec.get("type", "")

    if line_width > 0:
        shape.line.color.rgb = RGBColor.from_string(spec.get("line_hex", "000000"))
        shape.line.width     = Pt(line_width)

        dash = spec.get("line_dash", "solid")
        if dash != "solid":
            spPr = shape.element.spPr
            ln   = spPr.find(qn("a:ln"))
            if ln is None:
                ln = etree.SubElement(spPr, qn("a:ln"))
            for old in ln.findall(qn("a:prstDash")):
                ln.remove(old)
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", dash)

    elif stype not in _CONNECTOR_TYPES:
        # Explicitly suppress the border so the shape renders without a stroke
        spPr = shape.element.spPr
        ln   = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        for child in list(ln):
            ln.remove(child)
        etree.SubElement(ln, qn("a:noFill"))


def _apply_text_runs(shape, spec: dict) -> None:
    text_runs = spec.get("text_runs")
    if not text_runs or not hasattr(shape, "text_frame"):
        return

    tf             = shape.text_frame
    tf.word_wrap   = True
    tf.auto_size   = None
    tf.vertical_anchor = _VALIGN_MAP.get(
        spec.get("v_align", "top"), MSO_ANCHOR.TOP
    )

    current_para  = tf.paragraphs[0]
    first_in_para = True

    for item in text_runs:
        if item.get("paragraph_break"):
            current_para  = tf.add_paragraph()
            first_in_para = True
            continue

        if first_in_para:
            current_para.alignment = _ALIGN_MAP.get(
                item.get("align", "left"), PP_ALIGN.LEFT
            )
            pPr = current_para._p.get_or_add_pPr()

            if "line_spacing_pt" in item:
                lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
                etree.SubElement(lnSpc, qn("a:spcPts")).set(
                    "val", str(int(item["line_spacing_pt"] * 100))
                )
            if "space_before_pt" in item:
                spcBef = etree.SubElement(pPr, qn("a:spcBef"))
                etree.SubElement(spcBef, qn("a:spcPts")).set(
                    "val", str(int(item["space_before_pt"] * 100))
                )
            if "space_after_pt" in item:
                spcAft = etree.SubElement(pPr, qn("a:spcAft"))
                etree.SubElement(spcAft, qn("a:spcPts")).set(
                    "val", str(int(item["space_after_pt"] * 100))
                )
            first_in_para = False

        run             = current_para.add_run()
        run.text        = item["text"]
        run.font.name   = item.get("font_name", "Calibri")
        run.font.size   = Pt(item.get("font_size_pt", 12.0))
        run.font.bold   = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.underline = item.get("underline", False)
        if "font_color_hex" in item:
            run.font.color.rgb = RGBColor.from_string(item["font_color_hex"])


def _set_corner_radius(shape, spec: dict) -> None:
    corner_emu = spec.get("corner_radius_emu", 0)
    if not corner_emu:
        return
    prstGeom = shape.element.spPr.find(qn("a:prstGeom"))
    if prstGeom is None:
        return
    av_lst = prstGeom.find(qn("a:avLst"))
    if av_lst is None:
        av_lst = etree.SubElement(prstGeom, qn("a:avLst"))
    av = av_lst.find(qn("a:gd"))
    if av is None:
        av = etree.SubElement(av_lst, qn("a:gd"))
    av.set("name", "adj")
    min_dim = min(spec["cx"], spec["cy"])
    if min_dim > 0:
        av.set("val", str(min(50000, int(corner_emu / min_dim * 100000))))


# ── public API ───────────────────────────────────────────────────────────────

def build_slide(
    slide,
    specs: list[dict],
    original_image_bytes: bytes,
    slide_cx: int,
    slide_cy: int,
) -> None:
    """
    Populate a python-pptx Slide with shapes from a list of ShapeSpec dicts.
    The original slide image is inserted as a hidden fallback at the bottom of
    the z-stack so that non-reconstructable content degrades gracefully.
    """
    # ── Step 1: hidden fallback PNG at absolute bottom ───────────────────────
    pic     = slide.shapes.add_picture(
        BytesIO(original_image_bytes),
        left=0, top=0, width=slide_cx, height=slide_cy,
    )
    spTree  = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)               # index 2 = after nvGrpSpPr + grpSpPr
    pic._element.nvPicPr.cNvPr.set("hidden", "1")
    pic._element.nvPicPr.cNvPr.set("descr",  "__fallback_png__")

    # ── Step 2: build shapes sorted by z_order ───────────────────────────────
    shape_id_map: dict[int, object] = {}
    deferred:     list[tuple[dict, object]] = []

    for spec in sorted(specs, key=lambda s: s.get("z_order", 0)):
        stype = spec.get("type", "")
        shape = None

        try:
            if stype in _AUTO_SHAPE_MAP:
                shape = slide.shapes.add_shape(
                    _AUTO_SHAPE_MAP[stype],
                    spec["x"], spec["y"], spec["cx"], spec["cy"],
                )
                if stype == "rounded_rect":
                    _set_corner_radius(shape, spec)

            elif stype == "textbox":
                shape = slide.shapes.add_textbox(
                    spec["x"], spec["y"], spec["cx"], spec["cy"],
                )

            elif stype in _CONNECTOR_TYPES:
                conn_type = _CONNECTOR_MAP.get(
                    spec.get("connector_type", "straight"),
                    MSO_CONNECTOR_TYPE.STRAIGHT,
                )
                shape = slide.shapes.add_connector(
                    conn_type,
                    spec.get("start_x", spec["x"]),
                    spec.get("start_y", spec["y"]),
                    spec.get("end_x",   spec["x"] + spec["cx"]),
                    spec.get("end_y",   spec["y"] + spec["cy"]),
                )

            else:
                log.warning("build_slide: unknown type %r (id=%s), skipping",
                            stype, spec.get("id"))
                continue

        except Exception:
            log.exception("build_slide: failed to add shape id=%s type=%s",
                          spec.get("id"), stype)
            continue

        _apply_rotation(shape, spec)

        if stype not in _CONNECTOR_TYPES:
            _apply_fill(shape, spec)

        _apply_line(shape, spec)
        _apply_text_runs(shape, spec)

        spec_id = spec.get("id")
        if spec_id is not None:
            shape_id_map[spec_id] = shape

        if stype in _CONNECTOR_TYPES and (
            spec.get("start_shape_id") or spec.get("end_shape_id")
        ):
            deferred.append((spec, shape))

    # ── Step 3: connect deferred connector endpoints ─────────────────────────
    for spec, connector in deferred:
        start_id = spec.get("start_shape_id")
        end_id   = spec.get("end_shape_id")

        if start_id and start_id in shape_id_map:
            try:
                connector.begin_connect(
                    shape_id_map[start_id],
                    spec.get("start_anchor", 0),
                )
            except Exception:
                log.debug("begin_connect failed for connector id=%s", spec.get("id"))

        if end_id and end_id in shape_id_map:
            try:
                connector.end_connect(
                    shape_id_map[end_id],
                    spec.get("end_anchor", 2),
                )
            except Exception:
                log.debug("end_connect failed for connector id=%s", spec.get("id"))
