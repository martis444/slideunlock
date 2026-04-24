import logging
import os
import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation

from .ai_reconstructor import reconstruct
from .classifier import classify_all
from .harvester import harvest
from .shape_builder import build_slide
from .ssim_gate import verify_and_nudge
from .ungrouper import flatten_groups
from .xml_surgery import strip_locks

log = logging.getLogger(__name__)


def _patch_zip(zip_path: str, patches: dict[str, bytes]) -> None:
    """Apply {entry: new_bytes} patches to a zip in-place (no-op if empty)."""
    if not patches:
        return
    tmp = zip_path + "._patch"
    with zipfile.ZipFile(zip_path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = patches.get(info.filename, zin.read(info.filename))
                zout.writestr(info, data)
    os.replace(tmp, zip_path)


def _clear_slide_shapes(sp_tree) -> None:
    """Remove all shape children from spTree, keeping the two fixed headers."""
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)


def process_pptx(input_bytes: bytes) -> bytes:
    """
    Full unlock pipeline.

    1. Strip protection locks from all slides/layouts/masters.
    2. Flatten shape groups into individual shapes.
    3. For flat-image slides (slides that are a single full-bleed image),
       AI-reconstruct the shapes and verify with SSIM.

    Returns the processed PPTX bytes.
    """
    tmp_dir = tempfile.mkdtemp(prefix="slideunlock_")
    try:
        work_path = str(Path(tmp_dir) / "work.pptx")
        Path(work_path).write_bytes(input_bytes)

        # ── Step 1: strip locks ───────────────────────────────────────────────
        lock_patches = strip_locks(work_path)
        _patch_zip(work_path, lock_patches)
        log.info("Stripped locks: %d entries changed", len(lock_patches))

        # ── Step 2: flatten groups ────────────────────────────────────────────
        group_patches = flatten_groups(work_path)
        _patch_zip(work_path, group_patches)
        log.info("Flattened groups: %d entries changed", len(group_patches))

        # ── Step 3: harvest style context ─────────────────────────────────────
        with zipfile.ZipFile(work_path) as zf:
            style_ctx = harvest(zf)

        slide_cx = style_ctx["slide_cx_emu"]
        slide_cy = style_ctx["slide_cy_emu"]

        # ── Step 4: classify slides ───────────────────────────────────────────
        reports     = classify_all(work_path, style_ctx)
        flat_slides = [r for r in reports if r["is_flat_image"]]
        log.info(
            "Classified %d slides; %d require AI reconstruction",
            len(reports), len(flat_slides),
        )

        if not flat_slides:
            return Path(work_path).read_bytes()

        # ── Step 5: AI-reconstruct flat image slides ──────────────────────────
        for report in flat_slides:
            slide_num  = report["slide_num"]   # 1-based
            slide_idx  = slide_num - 1
            media_path = report.get("flat_image_media_path")

            if not media_path:
                log.warning("Slide %d: no embedded image path, skipping", slide_num)
                continue

            try:
                with zipfile.ZipFile(work_path) as zf:
                    image_bytes = zf.read(media_path)
            except KeyError:
                log.warning("Slide %d: image %r not found in zip", slide_num, media_path)
                continue

            log.info("Slide %d: calling AI reconstructor (image=%s)", slide_num, media_path)
            specs = reconstruct(image_bytes, style_ctx, slide_cx, slide_cy)
            if not specs:
                log.warning("Slide %d: AI returned no shapes, leaving as-is", slide_num)
                continue

            # Reload prs fresh from current zip state on each iteration so that
            # changes written by verify_and_nudge for prior slides are included.
            prs     = Presentation(work_path)
            slide   = prs.slides[slide_idx]
            sp_tree = slide.shapes._spTree

            _clear_slide_shapes(sp_tree)
            build_slide(slide, specs, image_bytes, slide_cx, slide_cy)

            # Save so verify_and_nudge can render via LibreOffice
            prs.save(work_path)

            ssim_ok, ssim_score, status = verify_and_nudge(
                original_slide_image=image_bytes,
                rebuilt_pptx_path=work_path,
                slide_index=slide_idx,
                slide_shapes_spTree=sp_tree,
                specs=specs,
            )
            log.info(
                "Slide %d: SSIM=%.4f  status=%s  passed=%s",
                slide_num, ssim_score, status, ssim_ok,
            )
            # verify_and_nudge has written the final slide XML to the zip.
            # The next iteration reloads prs from the updated zip.

        return Path(work_path).read_bytes()

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
