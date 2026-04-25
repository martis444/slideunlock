import logging
import shutil
import tempfile
import time
import zipfile

from .harvester import harvest
from .xml_surgery import strip_locks
from .ungrouper import flatten_groups
from .classifier import classify_all
from .ai_reconstructor import reconstruct
from .shape_builder import build_slide
from .ssim_gate import verify_and_nudge
from .repacker import repack

log = logging.getLogger(__name__)


def _clear_slide_shapes(sp_tree) -> None:
    """Remove all shape children from spTree, keeping the two fixed headers."""
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)


def unlock(
    input_path: str,
    output_path: str,
    basic_only: bool = True,
    reconstruct_flat: bool = False,
    ssim_threshold: float = 0.995,
) -> dict:
    start = time.time()
    tmp_dir = tempfile.mkdtemp()
    working_path = shutil.copy(input_path, tmp_dir + "/working.pptx")

    try:
        # Phase 0: Harvest
        with zipfile.ZipFile(working_path) as z:
            style_ctx = harvest(z)

        # Phase 1: Strip locks
        lock_changes = strip_locks(working_path)

        # Phase 2: Ungroup
        group_changes = flatten_groups(working_path)

        # Apply Phase 1+2 changes via surgical repack
        all_xml_changes = {**lock_changes, **group_changes}
        if all_xml_changes:
            repack(working_path, working_path + ".tmp", all_xml_changes)
            shutil.move(working_path + ".tmp", working_path)

        # Phase 3: Classify
        reports = classify_all(working_path, style_ctx)

        if basic_only or not reconstruct_flat:
            shutil.copy(working_path, output_path)
            return {
                "slides": [
                    {**r, "reconstruction_status": "skipped", "ssim_score": None}
                    for r in reports
                ],
                "locked_removed": sum(len(v) for v in lock_changes.values()),
                "groups_flattened": len(group_changes),
                "processing_time_seconds": time.time() - start,
            }

        # Phase 4+5+5.5: AI reconstruction for flat slides
        from pptx import Presentation

        prs = Presentation(working_path)
        slide_results = []
        # specs_per_slide stores the final specs list for each flat slide index
        specs_per_slide: dict[int, list] = {}

        for report in reports:
            idx = report["slide_num"] - 1
            slide = prs.slides[idx]

            if not report["is_flat_image"]:
                slide_results.append(
                    {**report, "reconstruction_status": "skipped", "ssim_score": None}
                )
                continue

            # Extract flat image bytes
            with zipfile.ZipFile(working_path) as z:
                img_bytes = z.read(report["flat_image_media_path"])

            # Phase 4: AI spec
            specs = reconstruct(
                img_bytes,
                style_ctx,
                style_ctx["slide_cx_emu"],
                style_ctx["slide_cy_emu"],
            )
            if not specs:
                log.warning("Slide %d: AI returned no shapes — leaving original image", report["slide_num"])
                slide_results.append(
                    {**report, "reconstruction_status": "fallback_png", "ssim_score": None}
                )
                continue

            # Phase 5: Build shapes into slide
            log.info("Slide %d: AI returned %d shapes — rebuilding", report["slide_num"], len(specs))
            _clear_slide_shapes(slide.shapes._spTree)
            build_slide(
                slide,
                specs,
                img_bytes,
                style_ctx["slide_cx_emu"],
                style_ctx["slide_cy_emu"],
            )
            specs_per_slide[idx] = specs
            slide_results.append(
                {**report, "reconstruction_status": "done", "ssim_score": None}
            )

        # Save after all Phase 5 builds
        rebuilt_path = tmp_dir + "/rebuilt.pptx"
        prs.save(rebuilt_path)

        # Phase 5.5: SSIM gate per flat slide
        for i, report in enumerate(reports):
            if not report["is_flat_image"]:
                continue
            idx = report["slide_num"] - 1
            if idx not in specs_per_slide:
                # Already marked fallback_png — nothing to verify
                continue

            with zipfile.ZipFile(working_path) as z:
                orig_img_bytes = z.read(report["flat_image_media_path"])

            passed, score, status = verify_and_nudge(
                original_slide_image=orig_img_bytes,
                rebuilt_pptx_path=rebuilt_path,
                slide_index=idx,
                slide_shapes_spTree=prs.slides[idx].shapes._spTree,
                specs=specs_per_slide[idx],
                ssim_threshold=ssim_threshold,
            )
            slide_results[i] = {**report, "reconstruction_status": status, "ssim_score": score}

        # Re-save after any nudge mutations
        prs.save(rebuilt_path)
        shutil.copy(rebuilt_path, output_path)

        return {
            "slides": slide_results,
            "locked_removed": sum(len(v) for v in lock_changes.values()),
            "groups_flattened": len(group_changes),
            "processing_time_seconds": time.time() - start,
        }

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    import json
    import sys

    result = unlock(
        sys.argv[1],
        sys.argv[2],
        basic_only="--full" not in sys.argv,
        reconstruct_flat="--full" in sys.argv,
    )
    print(json.dumps(result, indent=2))
